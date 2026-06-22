"""
Document Analysis Service - Phase 2
Handles AI-powered analysis of email attachments with tiered access
"""
import asyncio
import json
import logging
from typing import Optional, Dict, Any, List
from datetime import datetime
from io import BytesIO

from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select

from app.models.document_models import EmailAttachment, DocumentAnalysis
from app.models.database import User
from app.services.llm_orchestration_service import LLMOrchestrationService
from app.core.security import logger

logger = logging.getLogger(__name__)

# Optional dependencies for document processing
try:
    import PyPDF2
    HAS_PYPDF2 = True
except ImportError:
    HAS_PYPDF2 = False
    logger.warning("⚠️ PyPDF2 not installed - PDF processing disabled")

try:
    from docx import Document as DocxDocument
    HAS_PYTHON_DOCX = True
except ImportError:
    HAS_PYTHON_DOCX = False
    logger.warning("⚠️ python-docx not installed - DOCX processing disabled")

try:
    from openpyxl import load_workbook
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False
    logger.warning("⚠️ openpyxl not installed - XLSX processing disabled")

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    logger.warning("⚠️ python-pptx not installed - PPTX processing disabled")


class DocumentAnalysisEngine:
    """AI-powered document analysis with tiered access"""
    
    def __init__(self):
        self.llm_service = LLMOrchestrationService()
        # Analysis prompts optimized for document understanding
        self.analysis_prompt_template = """Analyze the following document and provide comprehensive insights.

Document: {filename}
Content Type: {file_type}
{content}

Please provide analysis in the following valid JSON format ONLY (no markdown, no extra text):
{{
  "summary": "2-3 sentence concise summary of the document",
  "key_points": ["point 1", "point 2", "point 3", "point 4", "point 5"],
  "entities": ["entity 1", "entity 2", "entity 3"],
  "sentiment": "positive|neutral|negative",
  "document_type": "email|report|contract|invoice|proposal|other",
  "confidence_score": 0.85,
  "language": "en"
}}"""

    async def analyze_document(
        self,
        session: AsyncSession,
        attachment: EmailAttachment,
        extracted_text: str,
        page_count: Optional[int],
        title: Optional[str],
        user_id: str,
        user_plan: str = "free"
    ) -> Optional[DocumentAnalysis]:
        """
        Perform AI analysis on document content
        
        Args:
            session: Database session
            attachment: EmailAttachment record
            extracted_text: Extracted document text
            page_count: Number of pages (if applicable)
            title: Extracted title
            user_id: User ID
            user_plan: User subscription plan (free/pro/enterprise)
        
        Returns:
            DocumentAnalysis record with results
        """
        try:
            # Create analysis record
            analysis = DocumentAnalysis(
                attachment_id=attachment.id,
                user_id=user_id,
                analysis_type="comprehensive",
                file_name=attachment.filename,
                file_extension=attachment.extension,
                file_size_display=self._format_file_size(attachment.file_size),
                extracted_title=title or attachment.filename,
                page_count=page_count,
                analysis_status="processing",
                is_full_analysis=False
            )
            
            session.add(analysis)
            await session.flush()
            
            logger.info(f"📊 Starting AI analysis for: {attachment.filename} (ID: {analysis.id})")
            
            # For free users, just return metadata
            if user_plan == "free":
                analysis.analysis_status = "completed"
                analysis.is_full_analysis = False
                logger.info(f"✅ Free tier analysis completed (metadata only): {analysis.id}")
                return analysis
            
            # For paid users, perform full AI analysis
            if not extracted_text or len(extracted_text.strip()) < 10:
                logger.warning(f"⚠️ Insufficient text for analysis: {attachment.filename}")
                analysis.analysis_status = "completed"
                analysis.is_full_analysis = False
                return analysis
            
            # Prepare analysis prompt
            prompt = self.analysis_prompt_template.format(
                filename=attachment.filename,
                file_type=attachment.extension.upper(),
                content=extracted_text[:5000]  # Limit to 5000 chars for token efficiency
            )
            
            # Call LLM with document analysis profile
            try:
                result = await self.llm_service.call_llm(
                    prompt=prompt,
                    system_prompt="You are a document analysis expert. Return ONLY valid JSON, no markdown formatting.",
                    temperature=0.3,  # Lower temperature for consistent analysis
                    max_tokens=800,
                    user_id=user_id,
                    feature="document_analysis",
                    session=session
                )
                
                if result.get("success"):
                    # Parse AI response
                    response_text = result.get("response", "{}").strip()
                    
                    # Handle markdown code blocks if present
                    if response_text.startswith("```json"):
                        response_text = response_text[7:]
                    if response_text.startswith("```"):
                        response_text = response_text[3:]
                    if response_text.endswith("```"):
                        response_text = response_text[:-3]
                    response_text = response_text.strip()
                    
                    try:
                        analysis_data = json.loads(response_text)
                        
                        # Store analysis results
                        analysis.summary = analysis_data.get("summary", "")[:1000]
                        analysis.key_points = analysis_data.get("key_points", [])[:5]
                        analysis.entities = analysis_data.get("entities", [])[:10]
                        analysis.sentiment = analysis_data.get("sentiment", "neutral").lower()
                        analysis.language = analysis_data.get("language", "en")
                        raw_confidence = analysis_data.get("confidence_score", 0.75)
                        try:
                            confidence_value = float(raw_confidence)
                        except (ValueError, TypeError):
                            confidence_value = 0.75
                        # Persist as integer percentage to match DB schema
                        if confidence_value <= 1:
                            analysis.confidence_score = int(round(confidence_value * 100))
                        else:
                            analysis.confidence_score = int(round(confidence_value))
                        analysis.is_full_analysis = True
                        analysis.analysis_status = "completed"
                        
                        logger.info(f"✅ AI analysis completed: {analysis.id}")
                        logger.debug(f"   Summary: {analysis.summary[:100]}...")
                        
                    except json.JSONDecodeError as e:
                        logger.warning(f"⚠️ Failed to parse AI response as JSON: {e}")
                        logger.debug(f"   Response was: {response_text[:200]}")
                        analysis.analysis_status = "error"
                        analysis.error_message = "Invalid AI response format"
                
                else:
                    analysis.analysis_status = "error"
                    analysis.error_message = result.get("error", "Unknown LLM error")
                    logger.error(f"❌ LLM call failed: {analysis.error_message}")
                
            except Exception as e:
                logger.error(f"❌ AI analysis exception: {e}")
                analysis.analysis_status = "error"
                analysis.error_message = str(e)[:500]
            
            return analysis
            
        except Exception as e:
            logger.error(f"❌ Error in analyze_document: {e}")
            return None
    
    @staticmethod
    def _format_file_size(bytes_size: int) -> str:
        """Format file size in human-readable format"""
        units = ['B', 'KB', 'MB', 'GB', 'TB']
        size = bytes_size
        unit_idx = 0
        
        while size >= 1024 and unit_idx < len(units) - 1:
            size /= 1024
            unit_idx += 1
        
        return f"{size:.1f} {units[unit_idx]}"


class DocumentTextExtractor:
    """Extract text from various document formats"""
    
    @staticmethod
    async def extract_pdf(file_content: bytes) -> tuple[str, Optional[int], Optional[str]]:
        """Extract text from PDF"""
        try:
            if not HAS_PYPDF2:
                logger.warning("⚠️ PyPDF2 not installed - PDF text extraction disabled")
                return "", None, None
            
            from io import BytesIO
            
            pdf_reader = PyPDF2.PdfReader(BytesIO(file_content))
            page_count = len(pdf_reader.pages)
            
            # Extract text from first page for title
            first_page_text = ""
            if page_count > 0:
                try:
                    first_page_text = pdf_reader.pages[0].extract_text() or ""
                except Exception as e:
                    logger.warning(f"Failed to extract first page: {e}")
            
            # Extract title (first non-empty line)
            title = None
            lines = first_page_text.split('\n')
            for line in lines:
                line = line.strip()
                if len(line) > 5 and len(line) < 200:
                    title = line
                    break
            
            # Extract full text
            full_text = ""
            for idx, page in enumerate(pdf_reader.pages[:20]):  # Limit to first 20 pages
                try:
                    page_text = page.extract_text() or ""
                    full_text += page_text + "\n"
                except Exception as e:
                    logger.debug(f"Failed to extract page {idx}: {e}")
            
            return full_text[:8000], page_count, title  # Limit to 8000 chars for analysis
            
        except Exception as e:
            logger.error(f"PDF extraction error: {e}")
            return "", None, None
    
    @staticmethod
    async def extract_docx(file_content: bytes) -> tuple[str, Optional[str]]:
        """Extract text from DOCX"""
        try:
            if not HAS_PYTHON_DOCX:
                logger.warning("⚠️ python-docx not installed - DOCX text extraction disabled")
                return "", None
            
            from io import BytesIO
            
            doc = DocxDocument(BytesIO(file_content))
            
            # Extract title from first paragraph
            title = None
            full_text = ""
            
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text += para.text + "\n"
                    if title is None and len(para.text.strip()) > 5:
                        title = para.text.strip()
            
            return full_text[:8000], title
            
        except Exception as e:
            logger.error(f"DOCX extraction error: {e}")
            return "", None
    
    @staticmethod
    async def extract_csv(file_content: bytes) -> str:
        """Extract contents from CSV"""
        try:
            text = file_content.decode('utf-8', errors='ignore')
            # Limit to first 1000 lines for analysis
            lines = text.split('\n')[:1000]
            return '\n'.join(lines)[:8000]
        except Exception as e:
            logger.error(f"CSV extraction error: {e}")
            return ""

    @staticmethod
    async def extract_xlsx(file_content: bytes) -> tuple[str, Optional[str]]:
        """Extract text-like values from XLSX workbook"""
        try:
            if not HAS_OPENPYXL:
                logger.warning("⚠️ openpyxl not installed - XLSX text extraction disabled")
                return "", None

            wb = load_workbook(filename=BytesIO(file_content), read_only=True, data_only=True)
            lines: List[str] = []
            title = None

            for sheet in wb.worksheets[:5]:
                lines.append(f"[Sheet] {sheet.title}")
                for row in sheet.iter_rows(min_row=1, max_row=200, values_only=True):
                    row_vals = [str(cell).strip() for cell in row if cell not in (None, "")]
                    if not row_vals:
                        continue
                    line = " | ".join(row_vals)
                    lines.append(line)
                    if title is None and len(line) > 5:
                        title = line[:120]
                    if len(lines) > 1200:
                        break
                if len(lines) > 1200:
                    break

            text = "\n".join(lines)[:8000]
            return text, title
        except Exception as e:
            logger.error(f"XLSX extraction error: {e}")
            return "", None

    @staticmethod
    async def extract_pptx(file_content: bytes) -> tuple[str, Optional[int], Optional[str]]:
        """Extract text from PPTX slides"""
        try:
            if not HAS_PPTX:
                logger.warning("⚠️ python-pptx not installed - PPTX text extraction disabled")
                return "", None, None

            prs = Presentation(BytesIO(file_content))
            lines: List[str] = []
            title = None
            slide_count = len(prs.slides)

            for idx, slide in enumerate(prs.slides[:50]):
                lines.append(f"[Slide {idx + 1}]")
                for shape in slide.shapes:
                    if not hasattr(shape, "text"):
                        continue
                    text = (shape.text or "").strip()
                    if not text:
                        continue
                    lines.append(text)
                    if title is None and len(text) > 5:
                        title = text[:120]
                if len(lines) > 2000:
                    break

            full_text = "\n".join(lines)[:8000]
            return full_text, slide_count, title
        except Exception as e:
            logger.error(f"PPTX extraction error: {e}")
            return "", None, None
    
    @staticmethod
    async def extract_text(file_content: bytes, file_ext: str) -> tuple[str, Optional[int], Optional[str]]:
        """
        Extract text from document based on file type
        
        Returns: (text, page_count, title)
        """
        file_ext = file_ext.lower()
        
        if file_ext == 'pdf':
            return await DocumentTextExtractor.extract_pdf(file_content)
        
        elif file_ext in ['docx', 'doc']:
            text, title = await DocumentTextExtractor.extract_docx(file_content)
            return text, None, title
        
        elif file_ext == 'csv':
            return await DocumentTextExtractor.extract_csv(file_content), None, None

        elif file_ext in ['xlsx', 'xlsm']:
            text, title = await DocumentTextExtractor.extract_xlsx(file_content)
            return text, None, title

        elif file_ext in ['pptx', 'ppt']:
            return await DocumentTextExtractor.extract_pptx(file_content)
        
        elif file_ext in ['txt', 'log']:
            try:
                text = file_content.decode('utf-8', errors='ignore')
                return text[:8000], None, None
            except Exception as e:
                logger.error(f"Text file extraction error: {e}")
                return "", None, None
        
        elif file_ext in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'webp']:
            # For images, return placeholder
            return "[Image file - visual inspection required]", None, None
        
        elif file_ext in ['zip', 'rar', '7z', 'gz']:
            return "[Archive file - extraction not yet supported]", None, None
        
        else:
            return f"[{file_ext.upper()} file - content extraction not supported]", None, None


class DocumentAnalysisBackgroundTask:
    """Background task handler for document analysis"""
    
    def __init__(self):
        self.text_extractor = DocumentTextExtractor()
        self.analysis_engine = DocumentAnalysisEngine()
    
    async def process_attachment_analysis(
        self,
        session: AsyncSession,
        attachment_id: str,
        user_id: str,
        user_plan: str = "free"
    ) -> Optional[DocumentAnalysis]:
        """
        Background task to analyze a single attachment
        
        Args:
            session: Database session
            attachment_id: EmailAttachment ID
            user_id: User ID
            user_plan: User subscription plan
        
        Returns:
            DocumentAnalysis record or None
        """
        try:
            # Get attachment
            stmt = select(EmailAttachment).where(EmailAttachment.id == attachment_id)
            result = await session.execute(stmt)
            attachment = result.scalars().first()
            
            if not attachment:
                logger.error(f"Attachment not found: {attachment_id}")
                return None
            
            # Check if analysis already exists
            stmt = select(DocumentAnalysis).where(
                DocumentAnalysis.attachment_id == attachment_id
            )
            result = await session.execute(stmt)
            existing = result.scalars().first()
            
            if existing:
                logger.debug(f"Analysis already exists for: {attachment_id}")
                return existing
            
            logger.info(f"🔄 Processing analysis for attachment: {attachment.filename}")
            
            # Read file
            from app.services.attachment_service import AttachmentService
            att_service = AttachmentService()
            file_content = att_service.read_file_content(attachment.storage_path)
            
            if not file_content:
                logger.error(f"Cannot read file for analysis: {attachment.filename}")
                return None
            
            # Extract text
            extracted_text, page_count, title = await self.text_extractor.extract_text(
                file_content,
                attachment.extension
            )
            
            logger.info(f"📝 Extracted {len(extracted_text)} characters from {attachment.filename}")
            
            # Perform AI analysis
            analysis = await self.analysis_engine.analyze_document(
                session,
                attachment,
                extracted_text,
                page_count,
                title,
                user_id,
                user_plan
            )
            
            if analysis:
                await session.commit()
                logger.info(f"✅ Analysis complete: {attachment.filename}")
            
            return analysis
            
        except Exception as e:
            logger.error(f"❌ Error in process_attachment_analysis: {e}")
            try:
                await session.rollback()
            except:
                pass
            return None
    
    async def process_email_attachments(
        self,
        session: AsyncSession,
        email_id: str,
        user_id: str,
        user_plan: str = "free"
    ) -> List[DocumentAnalysis]:
        """
        Background task to analyze all attachments for an email
        
        Args:
            session: Database session
            email_id: Email ID
            user_id: User ID
            user_plan: User subscription plan
        
        Returns:
            List of DocumentAnalysis records
        """
        try:
            # Get all attachments for email
            stmt = select(EmailAttachment).where(EmailAttachment.email_id == email_id)
            result = await session.execute(stmt)
            attachments = result.scalars().all()
            
            if not attachments:
                return []
            
            logger.info(f"📧 Processing {len(attachments)} attachments for email: {email_id}")
            
            analyses = []
            for attachment in attachments:
                try:
                    analysis = await self.process_attachment_analysis(
                        session,
                        attachment.id,
                        user_id,
                        user_plan
                    )
                    if analysis:
                        analyses.append(analysis)
                except Exception as e:
                    logger.error(f"Failed to process attachment {attachment.id}: {e}")
            
            logger.info(f"✅ Completed analysis for {len(analyses)} attachments")
            return analyses
            
        except Exception as e:
            logger.error(f"❌ Error in process_email_attachments: {e}")
            return []
